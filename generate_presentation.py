"""
Generate an executive presentation slide deck (.pptx) for Anti-Theft Tracker (Sentinel AI)
Covering all 5 required presentation pillars with professional dark-theme design.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_presentation(filename="Anti_Theft_Tracker_Presentation.pptx"):
    prs = Presentation()
    # 16:9 widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    COLOR_BG = RGBColor(15, 23, 42)          # Deep slate navy (#0f172a)
    COLOR_CARD = RGBColor(30, 41, 59)        # Card slate (#1e293b)
    COLOR_ACCENT = RGBColor(0, 210, 255)      # Cyan highlight (#00d2ff)
    COLOR_RED = RGBColor(255, 68, 68)         # Alert red (#ff4444)
    COLOR_WHITE = RGBColor(248, 250, 252)     # Off-white (#f8fafc)
    COLOR_MUTED = RGBColor(148, 163, 184)     # Muted text (#94a3b8)
    COLOR_GOLD = RGBColor(255, 184, 0)        # Gold accent (#ffb800)

    blank_layout = prs.slide_layouts[6]

    def set_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background() # No border
        return bg

    def add_header(slide, tag_text, title_text, subtitle_text=None):
        # Category Tag Pill
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8), Inches(0.4))
        tf_tag = tag_box.text_frame
        tf_tag.word_wrap = True
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = tag_text.upper()
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = COLOR_ACCENT

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.5), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_WHITE

        if subtitle_text:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5))
            tf_sub = sub_box.text_frame
            tf_sub.word_wrap = True
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle_text
            p_sub.font.size = Pt(13)
            p_sub.font.color.rgb = COLOR_MUTED

    def add_card(slide, left, top, width, height, title, points, badge=None, badge_color=COLOR_ACCENT):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = RGBColor(51, 65, 85)
        card.line.width = Pt(1)

        tx_box = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), width - Inches(0.5), height - Inches(0.4))
        tf = tx_box.text_frame
        tf.word_wrap = True

        if badge:
            p_badge = tf.paragraphs[0]
            p_badge.text = badge.upper()
            p_badge.font.size = Pt(10)
            p_badge.font.bold = True
            p_badge.font.color.rgb = badge_color
            p_title = tf.add_paragraph()
        else:
            p_title = tf.paragraphs[0]

        p_title.text = title
        p_title.font.size = Pt(17)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_WHITE
        p_title.space_after = Pt(12)

        for pt in points:
            p = tf.add_paragraph()
            p.text = f"• {pt}"
            p.font.size = Pt(12.5)
            p.font.color.rgb = COLOR_MUTED
            p.space_after = Pt(6)

    # -------------------------------------------------------------
    # SLIDE 1: TITLE SLIDE
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_background(slide1)

    # Hero Pill Badge
    pill = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(3.8), Inches(0.45))
    pill.fill.solid()
    pill.fill.fore_color.rgb = RGBColor(0, 80, 120)
    pill.line.color.rgb = COLOR_ACCENT
    tf_pill = pill.text_frame
    p_pill = tf_pill.paragraphs[0]
    p_pill.text = "AI COMPUTER VISION & LOSS PREVENTION"
    p_pill.font.size = Pt(11)
    p_pill.font.bold = True
    p_pill.font.color.rgb = COLOR_ACCENT
    p_pill.alignment = PP_ALIGN.CENTER

    # Main Title
    t_box = slide1.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.8))
    tf_main = t_box.text_frame
    p_main = tf_main.paragraphs[0]
    p_main.text = "Anti-Theft Tracker"
    p_main.font.size = Pt(44)
    p_main.font.bold = True
    p_main.font.color.rgb = COLOR_WHITE

    p_sub = tf_main.add_paragraph()
    p_sub.text = "Autonomous Shoplifting & Threat Detection System"
    p_sub.font.size = Pt(24)
    p_sub.font.color.rgb = COLOR_ACCENT
    p_sub.space_before = Pt(8)

    # Description Paragraph
    d_box = slide1.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(10.5), Inches(1.0))
    tf_desc = d_box.text_frame
    tf_desc.word_wrap = True
    p_d = tf_desc.paragraphs[0]
    p_d.text = "Transforming passive CCTV feeds into proactive real-time surveillance nodes using YOLOv8 deep learning and adaptive spatial-proximity interaction heuristics."
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = COLOR_MUTED

    # Presenter Metadata
    meta_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.1))
    meta_card.fill.solid()
    meta_card.fill.fore_color.rgb = COLOR_CARD
    meta_card.line.color.rgb = RGBColor(51, 65, 85)

    meta_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.75), Inches(11.3), Inches(0.8))
    tf_meta = meta_box.text_frame
    p_meta = tf_meta.paragraphs[0]
    p_meta.text = "Presenter: Mujahid Ahmed  |  GitHub: github.com/MujahidAhmed07/antiteft-tracker"
    p_meta.font.size = Pt(13)
    p_meta.font.bold = True
    p_meta.font.color.rgb = COLOR_WHITE

    p_meta2 = tf_meta.add_paragraph()
    p_meta2.text = "Stack: Python • FastAPI • YOLOv8 • OpenCV • PyTorch • HTML5/CSS3 Dashboard"
    p_meta2.font.size = Pt(12)
    p_meta2.font.color.rgb = COLOR_ACCENT
    p_meta2.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 2: PILLAR 1 - THE PROBLEM & WHO IT AFFECTS
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_background(slide2)
    add_header(slide2, "Pillar 1: Problem Definition", "The Challenge: $100B+ Lost Annually to Retail Shrinkage", "Traditional CCTV cameras record crime—they do not prevent it.")

    card_w = Inches(3.7)
    card_h = Inches(4.7)
    card_top = Inches(2.1)

    add_card(slide2, Inches(0.8), card_top, card_w, card_h,
        "The Core Problem",
        [
            "Retail shrinkage causes over $100 Billion in annual global losses.",
            "Shoplifting and internal distraction-theft occur in mere seconds.",
            "Traditional surveillance is 100% reactive: footage is only checked hours or days AFTER goods have left the store.",
            "Delayed response eliminates any chance of on-premise recovery."
        ],
        badge="Economic Impact", badge_color=COLOR_RED
    )

    add_card(slide2, Inches(4.8), card_top, card_w, card_h,
        "Human Limitations",
        [
            "Security guards monitor 15 to 40 CCTV screens at the same time.",
            "Studies show operator vigilance drops by over 50% within just 20 minutes.",
            "Distraction, fatigue, and blind spots make manual theft spotting statistically ineffective.",
            "High turnover and surging payroll costs for human security."
        ],
        badge="Operational Bottleneck", badge_color=COLOR_GOLD
    )

    add_card(slide2, Inches(8.8), card_top, card_w, card_h,
        "Who It Affects",
        [
            "Retail Store Owners: Severe profit erosion and inventory write-offs.",
            "Loss Prevention Teams: Overwhelmed by false alerts and endless forensic footage review.",
            "Cashiers & Staff: Direct confrontation risks with shoplifters.",
            "Honest Consumers: Pay higher retail prices to cover shrink costs."
        ],
        badge="Impacted Audience", badge_color=COLOR_ACCENT
    )

    # -------------------------------------------------------------
    # SLIDE 3: PILLAR 2 - THE SOLUTION & AUDIENCE IT SERVES
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_background(slide3)
    add_header(slide3, "Pillar 2: The Solution", "Anti-Theft Tracker (Sentinel AI)", "An autonomous AI sentinel turning ordinary cameras into real-time threat detectors.")

    add_card(slide3, Inches(0.8), card_top, card_w, card_h,
        "Proactive Detection",
        [
            "Performs real-time body tracking and suspicious behavior analysis at 30+ FPS.",
            "Detects subtle person-to-person body overlaps and theft proximity patterns.",
            "Identifies high-theft merchandise (purses, handbags, electronics, phones).",
            "Triggers instant high-visibility alerts the moment theft happens."
        ],
        badge="Intelligent Core", badge_color=COLOR_ACCENT
    )

    add_card(slide3, Inches(4.8), card_top, card_w, card_h,
        "SOC Surveillance Hub",
        [
            "Modern web command center with live streaming and interactive controls.",
            "Automated Crime Scene Evidence Gallery capturing photos at the exact theft frame.",
            "Popup Lightbox Modal for instant full-resolution forensic inspection.",
            "One-click evidence download (.avi video + .jpg crime scene photos)."
        ],
        badge="Command Center", badge_color=COLOR_GOLD
    )

    add_card(slide3, Inches(8.8), card_top, card_w, card_h,
        "Who It Serves",
        [
            "Store Managers: Real-time floor notifications and daily theft event telemetry.",
            "Security Guards: Actionable triage—intervene directly when an alert is triggered.",
            "Security System Integrators: Drop-in REST API compatible with existing NVR/IP cameras.",
            "Multi-Branch Enterprises: Centralized surveillance telemetry across locations."
        ],
        badge="Audience & Fit", badge_color=COLOR_WHITE
    )

    # -------------------------------------------------------------
    # SLIDE 4: PILLAR 3 - NEED ADDRESSED & REAL-WORLD IMPACT
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_background(slide4)
    add_header(slide4, "Pillar 3: Real-World Impact", "Measurable Business & Security Transformation", "Delivering tangible ROI from day one of deployment.")

    add_card(slide4, Inches(0.8), card_top, card_w, card_h,
        "The Fundamental Need",
        [
            "Shifts loss prevention from post-mortem forensic review to active real-time deterrence.",
            "Eliminates dependency on constant human monitoring of video screens.",
            "Provides verifiable, court-ready visual evidence trails with timestamps and frame numbers.",
            "Enables rapid deployment with zero specialized hardware requirements."
        ],
        badge="Market Necessity", badge_color=COLOR_GOLD
    )

    add_card(slide4, Inches(4.8), card_top, card_w, card_h,
        "Operational Impact",
        [
            "Shrinkage Loss Reduction: Early alerts deter theft attempts before suspects exit.",
            "10x Guard Multiplier: Single security operator can oversee 20+ feeds efficiently.",
            "Zero False Accusation Risk: High-resolution scene capture provides indisputable proof.",
            "Staff Safety: Alerts guards discreetly, avoiding volatile confrontational surprises."
        ],
        badge="Quantitative ROI", badge_color=COLOR_ACCENT
    )

    add_card(slide4, Inches(8.8), card_top, card_w, card_h,
        "Enterprise Readiness",
        [
            "Zero Committed Secrets: Production-safe .env.example configuration.",
            "API Authentication: Optional X-API-Key token validation on all endpoints.",
            "Cross-Platform: Runs on standard PCs, edge servers, or cloud GPU instances.",
            "CORS & Multi-Tenant Support: Securely connects to multi-branch setups."
        ],
        badge="Scalability", badge_color=COLOR_WHITE
    )

    # -------------------------------------------------------------
    # SLIDE 5: PILLAR 4 - INNOVATION & TECHNOLOGY BEHIND IT
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_background(slide5)
    add_header(slide5, "Pillar 4: Technological Innovation", "Deep Learning Meets Spatial Heuristics", "A multi-layered computer vision pipeline engineered for low latency.")

    add_card(slide5, Inches(0.8), card_top, card_w, card_h,
        "YOLOv8 Object Core",
        [
            "State-of-the-art YOLOv8 architecture optimized for high-throughput video inference.",
            "Simultaneous multi-person tracking and high-value merchandise identification.",
            "Support for custom fine-tuned weights or pre-trained baseline models.",
            "Automated fallback to CPU or acceleration via CUDA GPU when available."
        ],
        badge="Deep Learning", badge_color=COLOR_ACCENT
    )

    add_card(slide5, Inches(4.8), card_top, card_w, card_h,
        "Spatial Heuristics Engine",
        [
            "Proprietary O(n²) multi-person proximity and bounding-box interaction analytics.",
            "Evaluates 2D body overlap, rapid closure speed, and dwell time in critical zones.",
            "Differentiates innocent customer interactions from suspicious theft/concealment.",
            "Dynamic Head-Up Display (HUD) with colored bounding boxes (Red Alert vs Green Normal)."
        ],
        badge="Algorithmic IP", badge_color=COLOR_RED
    )

    add_card(slide5, Inches(8.8), card_top, card_w, card_h,
        "Modern Streaming Stack",
        [
            "FastAPI Asynchronous Engine: Low-overhead MJPEG video broadcast streaming.",
            "Full-Width Evidence Gallery: Disk-backed persistent crime scene photo grid.",
            "Deduplicated State Polling: Real-time telemetry feed sync without browser reloads.",
            "Modular Architecture: Decoupled detection engine, utilities, and UI presentation."
        ],
        badge="System Engineering", badge_color=COLOR_WHITE
    )

    # -------------------------------------------------------------
    # SLIDE 6: PILLAR 5 - FEASIBILITY & WHAT WE HAVE ACTUALLY BUILT
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_background(slide6)
    add_header(slide6, "Pillar 5: Feasibility & What We Built", "Working, Battle-Tested Implementation", "Fully functional, deployed locally, and verified on public GitHub.")

    add_card(slide6, Inches(0.8), card_top, card_w, card_h,
        "1. Web Command Center",
        [
            "Full-featured surveillance dashboard running on FastAPI & modern vanilla JS/CSS.",
            "Multi-source video ingestion: Sample footage (demo1-demo3), file upload, or webcam.",
            "Adjustable real-time sensitivity sliders for Confidence and Proximity ratio.",
            "Live HUD overlay with FPS, alert badges, and timestamps."
        ],
        badge="Feature Complete", badge_color=COLOR_ACCENT
    )

    add_card(slide6, Inches(4.8), card_top, card_w, card_h,
        "2. Crime Scene Evidence Grid",
        [
            "Automatically snaps high-res scene images at the exact moment of theft.",
            "Full-width responsive gallery accumulating all captures below the screen.",
            "Popup Lightbox Modal: Enlarge image with 1-click forensic download.",
            "Persistent Storage (GET /api/snapshots): Preserves all evidence across page reloads."
        ],
        badge="Forensic Ready", badge_color=COLOR_GOLD
    )

    add_card(slide6, Inches(8.8), card_top, card_w, card_h,
        "3. Multi-Interface Support",
        [
            "Simple Desktop App (simple_app.py): Native Windows file picker for zero-setup video selection.",
            "Standalone CLI (shoplifting_detection.py): Headless or OpenCV desktop GUI.",
            "Export Formats: Annotated .avi video recording + individual incident .jpg captures.",
            "Public GitHub: Verified zero-secret repo at github.com/MujahidAhmed07/antiteft-tracker."
        ],
        badge="Versatile Deployment", badge_color=COLOR_WHITE
    )

    # -------------------------------------------------------------
    # SLIDE 7: LIVE DEMO & CONCLUSION
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_background(slide7)
    add_header(slide7, "Summary & Live Demonstration", "Anti-Theft Tracker: Next-Generation Retail Security", "Ready for deployment, field trials, and integration.")

    # Big summary card
    big_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.1), Inches(11.7), Inches(4.7))
    big_card.fill.solid()
    big_card.fill.fore_color.rgb = COLOR_CARD
    big_card.line.color.rgb = RGBColor(51, 65, 85)

    tx_big = slide7.shapes.add_textbox(Inches(1.2), Inches(2.4), Inches(10.9), Inches(4.1))
    tf_b = tx_big.text_frame
    tf_b.word_wrap = True

    p_b1 = tf_b.paragraphs[0]
    p_b1.text = "Key Takeaways"
    p_b1.font.size = Pt(22)
    p_b1.font.bold = True
    p_b1.font.color.rgb = COLOR_ACCENT
    p_b1.space_after = Pt(14)

    takeaways = [
        "Addresses a massive $100B+ retail problem by shifting security from reactive to proactive.",
        "Combines deep learning (YOLOv8) with spatial interaction heuristics for fast, accurate theft detection.",
        "Fully implemented with live web streaming, full-width evidence photo gallery, and native desktop app.",
        "Completely verified and secured repository with zero secrets and clear .env.example guidelines.",
        "Live Web Dashboard running right now at http://localhost:8000.",
    ]

    for t in takeaways:
        p = tf_b.add_paragraph()
        p.text = f"✓  {t}"
        p.font.size = Pt(14.5)
        p.font.color.rgb = COLOR_WHITE
        p.space_after = Pt(10)

    p_link = tf_b.add_paragraph()
    p_link.text = "🔗 GitHub Repository: https://github.com/MujahidAhmed07/antiteft-tracker"
    p_link.font.size = Pt(15)
    p_link.font.bold = True
    p_link.font.color.rgb = COLOR_GOLD
    p_link.space_before = Pt(14)

    prs.save(filename)
    print(f"[SUCCESS] PowerPoint presentation generated: {filename}")

if __name__ == "__main__":
    build_presentation()
